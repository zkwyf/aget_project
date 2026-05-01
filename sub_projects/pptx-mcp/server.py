# pptx_server.py
import io
import os
import subprocess
import sys # Import sys for stderr printing
import tempfile
import shutil # For finding soffice
import base64 # <-- Add this import
from pathlib import Path
from typing import Optional, List, Union, Tuple, Dict, Any

import pptx # Import the module directly
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.dml.color import RGBColor # <-- Add RGBColor
from pptx.shapes.base import BaseShape

from fastmcp import FastMCP, Image, Context
from fastmcp.resources import FileResource
from dotenv import load_dotenv
from starlette.responses import JSONResponse # 增加为了自定义路由返回

# --- Configuration ---
# Directory to store generated presentations
# SAVE_DIR = Path("./presentations")
# SAVE_DIR.mkdir(exist_ok=True)
load_dotenv()

# --- !! IMPORTANT !! LibreOffice Dependency ---
# Set this if 'soffice' is not in your system PATH
SOFFICE_PATH = None # e.g., "/usr/bin/soffice" or "C:\\Program Files\\LibreOffice\\program\\soffice.bin"
# --- End Configuration ---


# --- FastMCP Server Definition ---
# NOTE: python-pptx is included, but LibreOffice CANNOT be installed via pip.
# Deployment requires manual installation of LibreOffice on the host system.
#
# Server host/port can be set via HOST and PORT environment variables, defaulting to 127.0.0.1:8000
HOST = os.environ.get("HOST", "127.0.0.1")
PORT = int(os.environ.get("PORT", "8000"))

mcp = FastMCP(
    "PowerPoint Creator 📊 (with Image Rendering)",
    dependencies=["python-pptx", "pillow"] # Added pillow dependency for Image helper
)

# --- Helper Functions ---

def _find_soffice() -> str:
    """Finds the LibreOffice executable."""
    if SOFFICE_PATH and Path(SOFFICE_PATH).exists():
        return SOFFICE_PATH
    
    soffice_cmd = "soffice"
    if os.name == 'nt': # Windows
        soffice_cmd = "soffice.exe" # Check PATH
        # Common Windows install locations (add more if needed)
        program_files = os.environ.get("ProgramFiles", "C:\\Program Files")
        program_files_x86 = os.environ.get("ProgramFiles(x86)", "C:\\Program Files (x86)")
        possible_paths = [
            Path(program_files) / "LibreOffice" / "program" / "soffice.exe",
            Path(program_files_x86) / "LibreOffice" / "program" / "soffice.exe",
        ]
        for p in possible_paths:
            if p.exists():
                return str(p)

    # For Linux/macOS, shutil.which checks the PATH
    found_path = shutil.which(soffice_cmd)
    if found_path:
        return found_path

    raise RuntimeError(
        "LibreOffice 'soffice' executable not found in PATH or configured SOFFICE_PATH. "
        "Image rendering requires LibreOffice installation."
    )


def _get_presentation_path(filename: str) -> Path:
    """Constructs the full path for the presentation file."""
    if not filename.endswith(".pptx"):
        filename += ".pptx"
    file_path = Path(filename)
    os.makedirs(file_path.parent, exist_ok=True)
    return file_path

def _load_presentation(filename: str) -> pptx.Presentation: # Use qualified name
    """Loads a presentation or creates a new one if it doesn't exist."""
    path = _get_presentation_path(filename)
    if path.exists():
        try:
            return pptx.Presentation(str(path)) # Use qualified name
        except pptx.exc.PackageNotFoundError: # Specific exception for file not found or not a zip file
            # This can happen if the file is not a valid pptx or is corrupted in a way that it's not recognized as a zip
            raise ValueError(f"File '{filename}' is not a valid PPTX file or is corrupted.")
        except Exception as e:
            # Catch other potential errors during loading (e.g., file corruption)
            raise ValueError(f"Error loading presentation '{filename}': {e}")
    else:
        return pptx.Presentation() # Use qualified name

def _save_presentation(prs: pptx.Presentation, filename: str): # Use qualified name
    """Saves the presentation object to the specified file."""
    path = _get_presentation_path(filename)
    try:
        prs.save(path)
    except Exception as e:
        raise IOError(f"Error saving presentation to '{path}': {e}")


def _get_slide(prs: pptx.Presentation, slide_index: int): # Use qualified name
    """Gets a specific slide by index, raising user-friendly errors."""
    if not isinstance(slide_index, int) or slide_index < 0:
        raise ValueError(f"Slide index must be a non-negative integer, got {slide_index}.")
    try:
        # Presentation.slides is list-like, check bounds explicitly
        if slide_index >= len(prs.slides):
             raise IndexError # Trigger the same error path
        return prs.slides[slide_index]
    except IndexError:
        raise ValueError(f"Invalid slide index {slide_index}. Presentation has {len(prs.slides)} slides (0-indexed).")

def _parse_shape_type(shape_name: str) -> MSO_SHAPE:
    """Converts a string shape name to an MSO_SHAPE enum."""
    try:
        # Convert to upper case for case-insensitive matching
        return getattr(MSO_SHAPE, shape_name.upper())
    except AttributeError:
        # Provide a list of common/useful shapes in the error
        common_shapes = ["RECTANGLE", "OVAL", "ROUNDED_RECTANGLE", "DIAMOND",
                         "ISOSCELES_TRIANGLE", "RIGHT_ARROW", "LEFT_ARROW",
                         "UP_ARROW", "DOWN_ARROW", "PENTAGON", "HEXAGON",
                         "CHEVRON", "STAR_5_POINT", "FLOWCHART_PROCESS",
                         "FLOWCHART_DECISION", "FLOWCHART_TERMINATOR",
                         "FLOWCHART_DATA", "LINE_CALLOUT_1"] # etc.
        raise ValueError(f"Unknown shape type '{shape_name}'. Try one of: {', '.join(common_shapes)}...")

def _get_shape_by_id(slide: pptx.slide.Slide, shape_id: int):
    """Finds a shape on the slide by its unique ID."""
    for shape in slide.shapes:
        if shape.shape_id == shape_id:
            return shape
    raise ValueError(f"Shape with ID {shape_id} not found on slide {slide.slide_id}. Available IDs: {[s.shape_id for s in slide.shapes]}")

# --- Helper functions for batch_update ---

def _resolve_slide_obj(
    prs: pptx.Presentation,
    page_ref: Union[str, int],
    object_map: Dict[str, Any],
    filename: str
) -> pptx.slide.Slide:
    """Resolves a page reference (string ID or int index) to a Slide object."""
    slide_obj: Optional[pptx.slide.Slide] = None
    if isinstance(page_ref, int):
        if 0 <= page_ref < len(prs.slides):
            slide_obj = prs.slides[page_ref]
        else:
            raise ValueError(f"Invalid slide index {page_ref} for presentation '{filename}'. It has {len(prs.slides)} slides.")
    elif isinstance(page_ref, str):
        resolved_item = object_map.get(page_ref)
        if isinstance(resolved_item, pptx.slide.Slide):
            slide_obj = resolved_item
        elif isinstance(resolved_item, int): # If we stored slide index by mistake for string ID
             if 0 <= resolved_item < len(prs.slides):
                slide_obj = prs.slides[resolved_item]
             else:
                raise ValueError(f"Resolved slide index {resolved_item} from ID '{page_ref}' is out of bounds for '{filename}'.")
        else:
            raise ValueError(f"Slide reference ID '{page_ref}' not found or not a slide in object_map for batch operation.")

    if slide_obj is None: # Should be caught by earlier checks, but as a safeguard
        raise ValueError(f"Could not resolve slide reference '{page_ref}'.")
    return slide_obj

def _resolve_shape_obj(
    slide: pptx.slide.Slide,
    shape_ref: Union[str, int],
    object_map: Dict[str, Any]
) -> BaseShape:
    """Resolves a shape reference (string ID or int ID) to a Shape object."""
    actual_shape_id: Optional[int] = None
    if isinstance(shape_ref, int):
        actual_shape_id = shape_ref
    elif isinstance(shape_ref, str):
        resolved_item = object_map.get(shape_ref)
        if isinstance(resolved_item, int): # Shape IDs are ints
            actual_shape_id = resolved_item
        elif isinstance(resolved_item, BaseShape): # If we stored the object itself
            actual_shape_id = resolved_item.shape_id
        else:
            raise ValueError(f"Shape reference ID '{shape_ref}' not found or not a shape ID in object_map for batch operation.")

    if actual_shape_id is None:
        raise ValueError(f"Could not determine actual shape ID from reference '{shape_ref}'.")

    return _get_shape_by_id(slide, actual_shape_id) # _get_shape_by_id will raise if not found on slide


# --- MCP Tools (Create, Add Slide, Add Elements - same as before) ---

@mcp.tool()
def create_or_clear_presentation(filename: str) -> str:
    """
    Creates a new, empty presentation with the given filename,
    or clears an existing one. Overwrites if the file exists.
    """
    # Ensure the presentation path is valid before creating
    pptx_path = _get_presentation_path(filename)
    prs = pptx.Presentation() # Use qualified name
    _save_presentation(prs, filename)
    return f"Presentation '{filename}' created/cleared."

@mcp.tool()
def add_slide(filename: str, layout_index: int = 5) -> str:
    """
    Adds a new slide to the presentation using a specified layout index.
    Common layouts: 0=Title, 1=Title+Content, 5=Title Only, 6=Blank.
    Returns the index of the newly added slide.
    """
    prs = _load_presentation(filename)
    if not (0 <= layout_index < len(prs.slide_layouts)):
        raise ValueError(f"Invalid layout_index {layout_index}. Must be between 0 and {len(prs.slide_layouts) - 1}.")
    slide_layout = prs.slide_layouts[layout_index]
    prs.slides.add_slide(slide_layout)
    new_slide_index = len(prs.slides) - 1 # Index is 0-based
    _save_presentation(prs, filename)
    return f"Added slide {new_slide_index} with layout {layout_index} to '{filename}'. New slide count: {len(prs.slides)}."


@mcp.tool()
def add_title_and_content(filename: str, slide_index: int, title: str, content: str) -> str:
    """
    Adds text to the title and main content placeholder of a specific slide.
    Assumes the slide layout has these placeholders (e.g., layout index 1).
    """
    prs = _load_presentation(filename)
    slide = _get_slide(prs, slide_index)

    title_shape = None
    content_placeholder = None

    # Find title placeholder (usually idx 0 or specific name)
    if slide.shapes.title:
        title_shape = slide.shapes.title
    else:
        for shape in slide.placeholders:
            if shape.name.lower().startswith("title"):
                title_shape = shape
                break

    if title_shape:
        title_shape.text = title
    else:
         print(f"Warning: Slide {slide_index} does not have a standard title placeholder.")


    # Find the main content placeholder (often index 1, but search robustly)
    for shape in slide.placeholders:
        # Check common indices or names
        if shape.placeholder_format.idx == 1 or \
           shape.name.lower().startswith("content placeholder") or \
           shape.name.lower().startswith("text placeholder") or \
           shape.name.lower().startswith("body"):

           # Avoid assigning the title shape if it was also found this way
           if shape != title_shape:
                content_placeholder = shape
                break

    if not content_placeholder and len(slide.placeholders) > 1 and slide.placeholders[1] != title_shape:
         # Fallback to index 1 if different from title
         content_placeholder = slide.placeholders[1]

    if content_placeholder:
        tf = content_placeholder.text_frame
        tf.text = content # Set first paragraph
        # Optionally clear other paragraphs if needed: while len(tf.paragraphs) > 1: tf._remove_paragraph(tf.paragraphs[-1])
        # Optionally add more paragraphs for bullet points if content has newlines etc.
    else:
        print(f"Warning: Slide {slide_index} does not seem to have a standard content placeholder.")
        # As a last resort, could add a new textbox, but maybe better to inform the user.

    _save_presentation(prs, filename)
    return f"Attempted to add title and content to slide {slide_index} in '{filename}'."


@mcp.tool()
def add_textbox(
    filename: str,
    slide_index: int,
    text: str,
    left_inches: float,
    top_inches: float,
    width_inches: float,
    height_inches: float,
    font_size_pt: int = 0, # Changed: Use 0 to indicate 'not set' instead of Optional[int]
    bold: bool = False
) -> str:
    """
    Adds a textbox with specified text, position, and dimensions (in inches) to a slide.
    Set font_size_pt to 0 or less to use the default font size.
    """ # Updated docstring
    prs = _load_presentation(filename)
    slide = _get_slide(prs, slide_index)
    left, top = Inches(left_inches), Inches(top_inches)
    width, height = Inches(width_inches), Inches(height_inches)

    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    # Handle multi-line text properly
    tf.text = text.split('\n')[0] # First line
    for line in text.split('\n')[1:]:
        p = tf.add_paragraph()
        p.text = line

    # Apply formatting to all paragraphs in the textbox
    for p in tf.paragraphs:
        p.font.bold = bold
        # Changed condition: Check for > 0 instead of checking for None
        if font_size_pt > 0:
            p.font.size = Pt(font_size_pt)

    _save_presentation(prs, filename)
    return f"Added textbox to slide {slide_index} in '{filename}'."


@mcp.tool()
def add_shape(
    filename: str,
    slide_index: int,
    shape_type_name: str,
    left_inches: float,
    top_inches: float,
    width_inches: float,
    height_inches: float,
    text: Optional[str] = None,
) -> str:
    """
    Adds an AutoShape (like RECTANGLE, OVAL, FLOWCHART_PROCESS) to a slide.
    Specify position and dimensions in inches. Optionally add text to the shape.
    Returns a confirmation message including the unique ID of the created shape.
    """ # Updated docstring
    prs = _load_presentation(filename)
    slide = _get_slide(prs, slide_index)
    shape_enum = _parse_shape_type(shape_type_name)
    left, top = Inches(left_inches), Inches(top_inches)
    width, height = Inches(width_inches), Inches(height_inches)

    shape = slide.shapes.add_shape(shape_enum, left, top, width, height)
    shape_id = shape.shape_id # Get the unique ID

    if text:
        tf = shape.text_frame
        # Handle multi-line text in shapes too
        tf.text = text.split('\n')[0]
        for line in text.split('\n')[1:]:
             p = tf.add_paragraph()
             p.text = line
        tf.word_wrap = True # Enable word wrap within the shape

    _save_presentation(prs, filename)
    # Include the shape_id in the return message
    return f"Added shape '{shape_type_name}' (ID: {shape_id}) to slide {slide_index} in '{filename}'."

@mcp.tool()
def add_connector(
    filename: str,
    slide_index: int,
    start_shape_id: int,
    end_shape_id: int,
    connector_type_name: str = "ELBOW", # Common default: STRAIGHT, ELBOW, CURVE
    start_connection_point_idx: int = 3, # Default: Mid-right side for many shapes
    end_connection_point_idx: int = 1,   # Default: Mid-left side for many shapes
) -> str:
    """
    Adds a connector shape between two existing shapes identified by their IDs.
    Defaults to an ELBOW connector from the right side of the start shape
    to the left side of the end shape.

    Args:
        filename: The presentation filename.
        slide_index: The 0-based index of the slide.
        start_shape_id: The unique ID of the shape where the connector starts.
        end_shape_id: The unique ID of the shape where the connector ends.
        connector_type_name: Type of connector (e.g., "STRAIGHT", "ELBOW", "CURVE").
        start_connection_point_idx: Index of the connection point on the start shape (0=center, 1-N=perimeter points).
        end_connection_point_idx: Index of the connection point on the end shape.

    Returns:
        Confirmation message including the connector's shape ID.
    """
    prs = _load_presentation(filename)
    slide = _get_slide(prs, slide_index)

    # Find the start and end shapes
    start_shape = _get_shape_by_id(slide, start_shape_id)
    end_shape = _get_shape_by_id(slide, end_shape_id)

    # Parse connector type
    try:
        connector_enum = getattr(MSO_CONNECTOR, connector_type_name.upper())
    except AttributeError:
        raise ValueError(f"Unknown connector type '{connector_type_name}'. Try: STRAIGHT, ELBOW, CURVE.")

    # Add the connector shape (initial position doesn't matter much)
    connector = slide.shapes.add_connector(
        connector_enum, Inches(1), Inches(1), Inches(2), Inches(2) # Arbitrary start/end
    )

    # Connect the shapes
    try:
        connector.begin_connect(start_shape, start_connection_point_idx)
    except Exception as e:
        # Provide more context on error
        print(f"Warning: Could not connect start of connector to shape {start_shape_id} at point {start_connection_point_idx}. Error: {e}. Ensure the connection point index is valid for the shape type.")
        # Attempt to connect to center as fallback?
        try:
            print("Attempting fallback connection to center (point 0) for start shape.")
            connector.begin_connect(start_shape, 0)
        except Exception as e2:
             print(f"Fallback connection to center also failed: {e2}")
             # Proceed without connection if fallback fails, user might fix later

    try:
        connector.end_connect(end_shape, end_connection_point_idx)
    except Exception as e:
        print(f"Warning: Could not connect end of connector to shape {end_shape_id} at point {end_connection_point_idx}. Error: {e}. Ensure the connection point index is valid for the shape type.")
        # Attempt to connect to center as fallback?
        try:
            print("Attempting fallback connection to center (point 0) for end shape.")
            connector.end_connect(end_shape, 0)
        except Exception as e2:
            print(f"Fallback connection to center also failed: {e2}")
            # Proceed without connection if fallback fails

    _save_presentation(prs, filename)
    connector_id = connector.shape_id
    return f"Added {connector_type_name} connector (ID: {connector_id}) from shape {start_shape_id} (point {start_connection_point_idx}) to shape {end_shape_id} (point {end_connection_point_idx}) on slide {slide_index}."

@mcp.tool()
def delete_shape(filename: str, slide_index: int, shape_id: int) -> str:
    """
    Deletes a specific shape from a slide using its unique ID.

    Args:
        filename: The presentation filename.
        slide_index: The 0-based index of the slide.
        shape_id: The unique ID of the shape to delete.

    Returns:
        Confirmation message.
    """
    prs = _load_presentation(filename)
    slide = _get_slide(prs, slide_index)

    # Find the shape to delete
    shape_to_delete = _get_shape_by_id(slide, shape_id)

    # Remove the shape element from its parent
    sp = shape_to_delete._element # Access the underlying XML element
    sp.getparent().remove(sp)   # Remove the element from the shapes collection

    _save_presentation(prs, filename)
    return f"Deleted shape with ID {shape_id} from slide {slide_index} in '{filename}'."

@mcp.tool()
def modify_shape(
    filename: str,
    slide_index: int,
    shape_id: int,
    text: Optional[str] = None,
    left_inches: Optional[float] = None,
    top_inches: Optional[float] = None,
    width_inches: Optional[float] = None,
    height_inches: Optional[float] = None,
    font_size_pt: Optional[int] = None,
    bold: Optional[bool] = None,
    fill_color_rgb: Optional[List[int]] = None,
    line_color_rgb: Optional[List[int]] = None,
    line_width_pt: Optional[float] = None
) -> str:
    """
    Modifies properties of an existing shape identified by its ID.
    Allows changing text, position, size, font attributes, fill color, and line style.
    Only provided parameters are changed.

    Args:
        filename: The presentation filename.
        slide_index: The 0-based index of the slide.
        shape_id: The unique ID of the shape to modify.
        text: New text content for the shape (replaces existing text).
        left_inches: New horizontal position from the left edge (in inches).
        top_inches: New vertical position from the top edge (in inches).
        width_inches: New width (in inches).
        height_inches: New height (in inches).
        font_size_pt: New font size (in points) for all text in the shape.
        bold: Set font bold state (True or False) for all text in the shape.
        fill_color_rgb: List of [R, G, B] values (0-255) for solid fill color.
        line_color_rgb: List of [R, G, B] values (0-255) for line color.
        line_width_pt: New line width (in points).

    Returns:
        Confirmation message summarizing the changes made.
    """
    prs = _load_presentation(filename)
    slide = _get_slide(prs, slide_index)
    shape = _get_shape_by_id(slide, shape_id)
    changes_made = []

    # Modify Position/Size
    if left_inches is not None:
        shape.left = Inches(left_inches)
        changes_made.append("position (left)")
    if top_inches is not None:
        shape.top = Inches(top_inches)
        changes_made.append("position (top)")
    if width_inches is not None:
        shape.width = Inches(width_inches)
        changes_made.append("size (width)")
    if height_inches is not None:
        shape.height = Inches(height_inches)
        changes_made.append("size (height)")

    # Modify Text and Font (if shape has text frame)
    if shape.has_text_frame:
        if text is not None:
            tf = shape.text_frame
            tf.clear() # Clear existing paragraphs before adding new text
            # Handle multi-line text properly
            lines = text.split('\n')
            tf.text = lines[0]
            for line in lines[1:]:
                p = tf.add_paragraph()
                p.text = line
            tf.word_wrap = True # Ensure word wrap is enabled
            changes_made.append("text content")

        font_changed = False
        if font_size_pt is not None:
            for p in shape.text_frame.paragraphs:
                p.font.size = Pt(font_size_pt)
            font_changed = True
        if bold is not None:
            for p in shape.text_frame.paragraphs:
                p.font.bold = bold
            font_changed = True
        if font_changed:
             changes_made.append("font attributes (size/bold)")

    # Modify Fill Color
    if fill_color_rgb is not None:
        if len(fill_color_rgb) == 3:
            fill = shape.fill
            fill.solid()
            fill.fore_color.rgb = RGBColor(*fill_color_rgb)
            changes_made.append("fill color")
        else:
            print(f"Warning: Invalid list {fill_color_rgb} for fill color. Expected [R, G, B] (length 3).")

    # Modify Line Style
    line_changed = False
    line = shape.line
    if line_color_rgb is not None:
        if len(line_color_rgb) == 3:
            line.color.rgb = RGBColor(*line_color_rgb)
            line_changed = True
        else:
             print(f"Warning: Invalid list {line_color_rgb} for line color. Expected [R, G, B] (length 3).")
    if line_width_pt is not None:
        line.width = Pt(line_width_pt)
        line_changed = True
    if line_changed:
        changes_made.append("line style (color/width)")

    if not changes_made:
        return f"No valid modifications specified for shape ID {shape_id} on slide {slide_index}."

    _save_presentation(prs, filename)
    return f"Modified shape ID {shape_id} on slide {slide_index}: updated {', '.join(changes_made)}."

# @mcp.tool()
# def add_picture(
#     filename: str,
#     slide_index: int,
#     image: bytes, # Changed type hint from Image to bytes
#     left_inches: float,
#     top_inches: float,
#     width_inches: Optional[float] = None,
#     height_inches: Optional[float] = None,
# ) -> str:
#     """
#     Adds a picture to a slide from provided image data.
#     Specify position in inches. Optionally specify width OR height in inches to scale.
#     If neither width nor height is given, the image's native size is used.
#     """
#     prs = _load_presentation(filename)
#     slide = _get_slide(prs, slide_index)
#     left, top = Inches(left_inches), Inches(top_inches)
#     width = Inches(width_inches) if width_inches is not None else None
#     height = Inches(height_inches) if height_inches is not None else None
#
#     # Use BytesIO to pass image data (which is now bytes) to python-pptx
#     image_stream = io.BytesIO(image) # Use image directly as it's now bytes
#
#     slide.shapes.add_picture(image_stream, left, top, width=width, height=height)
#
#     _save_presentation(prs, filename)
#     return f"Added picture to slide {slide_index} in '{filename}'."


@mcp.tool()
def add_picture(
    filename: str,
    slide_index: int,
    image_path: str,  # 改为文件路径
    left_inches: float,
    top_inches: float,
    width_inches: Optional[float] = None,
    height_inches: Optional[float] = None,
) -> str:
    """
    Adds a picture to a slide from an image file.
    Specify position in inches. Optionally specify width OR height in inches to scale.
    If neither width nor height is given, the image's native size is used.
    """
    prs = _load_presentation(filename)
    slide = _get_slide(prs, slide_index)
    left, top = Inches(left_inches), Inches(top_inches)
    width = Inches(width_inches) if width_inches is not None else None
    height = Inches(height_inches) if height_inches is not None else None

    # 直接从文件路径添加图片
    slide.shapes.add_picture(image_path, left, top, width=width, height=height)

    _save_presentation(prs, filename)
    return f"Added picture from '{image_path}' to slide {slide_index} in '{filename}'."


# --- MCP Resources (Description + NEW Image Rendering) ---

@mcp.resource("pptx://{filename}/slide/{slide_index}/description")
async def get_slide_content_description(filename: str, slide_index: str) -> str:
    """
    Provides a textual description of the shapes and text on a specific slide.
    Useful for the AI to 'confirm' the slide content without actual image rendering.
    [Reliable & Lightweight]
    """
    try:
        idx = int(slide_index)
    except ValueError:
        raise ValueError("Slide index must be an integer.")

    prs = _load_presentation(filename)
    slide = _get_slide(prs, idx) # _get_slide handles index errors

    description = f"--- Slide {idx} Content Description for '{filename}' ---\n"
    description += f"Layout: {slide.slide_layout.name}\n"
    description += f"Number of Shapes: {len(slide.shapes)}\n\n"

    for i, shape in enumerate(slide.shapes):
        shape_type = shape.shape_type
        # Use name attribute if available (e.g., for MSO_SHAPE enums), otherwise use string representation
        type_name = getattr(shape_type, 'name', str(shape_type))
        desc = f"Shape {i}: Type={type_name}"
        # Include the unique shape ID
        desc += f", ID={shape.shape_id}"
        try:
             desc += f", Left={shape.left.inches:.2f}\", Top={shape.top.inches:.2f}\", Width={shape.width.inches:.2f}\", Height={shape.height.inches:.2f}\""
        except AttributeError:
             desc += " (Position/Size not available)" # Handle shapes without these properties if they exist

        if shape.has_text_frame and shape.text.strip():
             # Truncate long text for brevity
            text_preview = (shape.text[:75] + '...') if len(shape.text) > 75 else shape.text
            desc += f", Text='{text_preview.replace(chr(11), ' ').replace('\n', ' ')}'" # Replace VT and newlines for single line desc

        description += desc + "\n"

    if not slide.shapes:
        description += "(Slide is empty)\n"

    description += "--- End Description ---"
    return description


@mcp.tool() # Changed from resource to tool
def get_slide_image(filename: str, slide_index: int) -> Image: # Changed slide_index type hint
    """
    Renders a specific slide as a PNG image using LibreOffice and returns it as an Image object.
    Requires LibreOffice installed and accessible on the server. May be slow.
    Use get_slide_content_description for a faster, text-based check.

    Args:
        filename: The name of the presentation file (e.g., "my_presentation.pptx").
        slide_index: The 0-based index of the slide to render.

    Returns:
        A fastmcp.Image object containing the PNG image data.

    Raises:
        ValueError: If filename is invalid, slide_index is not an integer or out of bounds.
        FileNotFoundError: If the presentation file or the generated PNG is not found.
        RuntimeError: If LibreOffice is not found or the conversion process fails/times out.
    """
    # Function body remains largely the same, just ensure idx is used correctly
    print(f"Attempting slide image rendering for slide {slide_index} of '{filename}' using LibreOffice...", file=sys.stderr)
    print("INFO: This requires LibreOffice to be installed and configured on the server.", file=sys.stderr)

    # Parameter is already slide_index: int, no need for conversion here
    idx = slide_index # Use the integer index directly

    pptx_path = _get_presentation_path(filename)
    if not pptx_path.exists():
        raise ValueError(f"Presentation file '{filename}' not found.")

    # Check number of slides BEFORE conversion to validate index early
    try:
        prs_check = pptx.Presentation(str(pptx_path)) # Use qualified name
        if idx >= len(prs_check.slides):
             raise ValueError(f"Invalid slide index {idx}. Presentation '{filename}' has {len(prs_check.slides)} slides (0-indexed).")
        del prs_check # Close file handle
    except Exception as e:
        print(f"ERROR: Failed to quickly check slide count for {filename}: {e}", file=sys.stderr)
        # Proceed cautiously, LibreOffice might handle corrupted files differently

    soffice = _find_soffice() # Raises RuntimeError if not found

    with tempfile.TemporaryDirectory() as temp_dir:
        temp_dir_path = Path(temp_dir)
        print(f"INFO: Using temporary directory: {temp_dir}", file=sys.stderr)

        # Command to convert the *entire* presentation to PNGs in the temp dir
        # soffice usually names output files based on the input filename + slide number
        cmd = [
            soffice,
            "--headless",          # Run without UI
            "--convert-to", "png", # Convert to PNG format
            "--outdir", str(temp_dir_path), # Output directory
            str(pptx_path)         # Input PPTX file
        ]

        try:
            print(f"INFO: Running LibreOffice command: {' '.join(cmd)}", file=sys.stderr)
            # Use a timeout to prevent hanging indefinitely
            timeout_seconds = 60
            process = subprocess.run(cmd, capture_output=True, text=True, check=False, timeout=timeout_seconds)

            if process.returncode != 0:
                print(f"ERROR: LibreOffice conversion failed! Return code: {process.returncode}", file=sys.stderr)
                print(f"ERROR: Stderr: {process.stderr}", file=sys.stderr)
                print(f"ERROR: Stdout: {process.stdout}", file=sys.stderr)
                raise RuntimeError(f"LibreOffice conversion failed (code {process.returncode}). Check MCP server logs for details.")
            else:
                 print("INFO: LibreOffice conversion process completed.", file=sys.stderr)
                 if process.stderr: # Often has warnings even on success
                     print(f"WARNING: LibreOffice stderr: {process.stderr}", file=sys.stderr)

        except FileNotFoundError:
            print(f"ERROR: '{soffice}' command not found. Ensure LibreOffice is installed and in PATH.", file=sys.stderr)
            raise RuntimeError("LibreOffice command failed: executable not found.")
        except subprocess.TimeoutExpired:
             print(f"ERROR: LibreOffice conversion timed out after {timeout_seconds} seconds.", file=sys.stderr)
             raise RuntimeError("LibreOffice conversion timed out.")
        except Exception as e:
             print(f"ERROR: An unexpected error occurred during LibreOffice execution: {e}", file=sys.stderr)
             raise RuntimeError(f"LibreOffice execution error: {e}")


        # Find the expected output file.
        # LibreOffice often names the output based on the input file name,
        # but may vary. A common pattern is just the input filename (without ext)
        # followed by the slide number (starting from 1 for the first slide!).
        # Or sometimes just the input filename if only one slide. Let's check robustly.
        base_filename = pptx_path.stem
        expected_png_filename = f"{base_filename}{idx + 1}.png" # soffice usually uses 1-based index for output
        expected_png_filename_single = f"{base_filename}.png" # Case for single-slide conversion output

        expected_png_path = temp_dir_path / expected_png_filename
        expected_png_path_single = temp_dir_path / expected_png_filename_single

        # List files to see what was actually created
        created_files = list(temp_dir_path.glob('*.png'))
        print(f"INFO: Files created in temp dir: {[f.name for f in created_files]}", file=sys.stderr)

        actual_png_path = None
        if expected_png_path.exists():
            actual_png_path = expected_png_path
        elif expected_png_path_single.exists() and len(created_files) == 1 and idx == 0:
            # If only one PNG was created and we asked for slide 0, assume it's the one
            actual_png_path = expected_png_path_single
        elif len(created_files) > idx:
             # Fallback: If soffice just numbered them sequentially without base name (less common)
             # Sort files to try and get consistent ordering (might be fragile)
             created_files.sort()
             potential_path = created_files[idx]
             print(f"WARNING: Could not find expected PNG file name, falling back to {potential_path.name} based on index.", file=sys.stderr)
             actual_png_path = potential_path
        elif created_files:
             # If some PNGs exist but not the one we expected
             print(f"WARNING: Expected PNG file {expected_png_filename} or {expected_png_filename_single} not found, but other PNGs exist.", file=sys.stderr)
             # Maybe take the first one if index is 0? Risky.
             if idx == 0:
                  actual_png_path = created_files[0]
                  print(f"WARNING: Using first found PNG: {actual_png_path.name}", file=sys.stderr)
             else:
                  raise FileNotFoundError(f"Could not determine the correct output PNG for slide {idx} in {temp_dir}. Found: {[f.name for f in created_files]}")
        else:
            raise FileNotFoundError(f"LibreOffice ran but no PNG output files were found in {temp_dir}.")


        print(f"INFO: Reading image data from: {actual_png_path}", file=sys.stderr)
        try:
            image_bytes = actual_png_path.read_bytes()
            # Return as FastMCP Image object
            return Image(data=image_bytes, format="png")
        except Exception as e:
             print(f"ERROR: Error reading PNG file {actual_png_path}: {e}", file=sys.stderr)
             raise RuntimeError(f"Failed to read generated PNG file: {e}")


@mcp.resource("pptx://{filename}/file", mime_type="application/vnd.openxmlformats-officedocument.presentationml.presentation")
async def get_pptx_file(filename: str):
    """
    Provides the .pptx file for download as a binary resource.
    """
    path = _get_presentation_path(filename).resolve()
    if not path.exists():
        raise FileNotFoundError(f"Presentation file '{filename}' not found.")
    return FileResource(
        uri=f"pptx://{filename}/file",
        path=path,
        is_binary=True,
        mime_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        name=f"PPTX file: {filename}",
        description=f"Download the PowerPoint file '{filename}'."
    )


@mcp.tool()
def get_presentation_file_b64(filename: str) -> str:
    """
    Returns the content of the presentation file encoded as a Base64 string.
    Useful if the client cannot handle MCP resources or direct file paths.

    Returns:
        A Base64 encoded string representing the binary content of the .pptx file.
    """
    path = _get_presentation_path(filename)
    if not path.exists():
        raise FileNotFoundError(f"Presentation file '{filename}' not found on the server.")

    try:
        file_bytes = path.read_bytes()
        base64_encoded_bytes = base64.b64encode(file_bytes)
        base64_string = base64_encoded_bytes.decode('utf-8') # Decode bytes to string for JSON compatibility
        return base64_string
    except Exception as e:
        raise IOError(f"Error reading or encoding file '{filename}': {e}")


# --- MCP Prompts (same as before) ---
@mcp.prompt()
def flowchart_guidance() -> str:
    """Provides tips on how to create flowcharts using the available tools."""
    return """
    To create a flowchart:
    1. Use `create_or_clear_presentation` to start.
    2. Use `add_slide` with a blank layout (e.g., layout_index=6).
    3. Use `add_shape` repeatedly for flowchart elements (e.g., FLOWCHART_PROCESS, FLOWCHART_DECISION, FLOWCHART_TERMINATOR). Use `text` parameter for labels. Position using `left_inches`, `top_inches`.
    4. Use `add_shape` with connector shapes (e.g., `RIGHT_ARROW`, or find specific connectors) OR use `add_connector` tool if available (not implemented here) to connect elements. Precise positioning takes care.
    5. Check arrangement using the description resource: `pptx://{filename}/slide/{slide_index}/description` (Reliable). # Keep description as resource
    6. Optionally, render the slide image using the tool: `get_slide_image` (Requires LibreOffice setup, may be slow). # Updated to tool name
    """

@mcp.prompt()
def available_shapes() -> str:
    """Lists some common MSO_SHAPE names usable with the add_shape tool."""
    common_shapes = [
        "RECTANGLE", "OVAL", "ROUNDED_RECTANGLE", "DIAMOND", "ISOSCELES_TRIANGLE",
        "RIGHT_ARROW", "LEFT_ARROW", "UP_ARROW", "DOWN_ARROW", "PENTAGON", "HEXAGON",
        "CHEVRON", "STAR_5_POINT", "LINE_CALLOUT_1", "ACTION_BUTTON_BACK_OR_PREVIOUS",
        # Flowchart shapes
        "FLOWCHART_PROCESS", "FLOWCHART_ALTERNATE_PROCESS", "FLOWCHART_DECISION",
        "FLOWCHART_DATA", "FLOWCHART_PREDEFINED_PROCESS", "FLOWCHART_INTERNAL_STORAGE",
        "FLOWCHART_DOCUMENT", "FLOWCHART_MULTIDOCUMENT", "FLOWCHART_TERMINATOR",
        "FLOWCHART_PREPARATION", "FLOWCHART_MANUAL_INPUT", "FLOWCHART_MANUAL_OPERATION",
        "FLOWCHART_CONNECTOR", "FLOWCHART_OFFPAGE_CONNECTOR", "FLOWCHART_CARD",
        "FLOWCHART_PUNCHED_TAPE", "FLOWCHART_SUMMING_JUNCTION", "FLOWCHART_OR",
        "FLOWCHART_COLLATE", "FLOWCHART_SORT", "FLOWCHART_EXTRACT", "FLOWCHART_MERGE",
        "FLOWCHART_STORED_DATA", "FLOWCHART_DELAY", "FLOWCHART_SEQUENTIAL_ACCESS_STORAGE",
        "FLOWCHART_MAGNETIC_DISK", "FLOWCHART_DIRECT_ACCESS_STORAGE", "FLOWCHART_DISPLAY"
    ]
    return f"Common shape names for `add_shape`: {', '.join(common_shapes)}. Many others exist."


# --- Batch Update Tool ---

@mcp.tool()
def batch_update(filename: str, requests: List[Dict[str, Dict[str, Any]]]) -> Dict[str, Any]:
    """
    Performs a series of operations on a presentation in a single batch.
    Similar to Google Slides API batchUpdate. Allows creating slides, adding shapes,
    text, pictures, connectors, modifying and deleting shapes.
    Operations are performed sequentially. If any operation fails,
    the entire batch is aborted, and no changes are saved to the presentation file.

    Args:
        filename: The name of the presentation file (e.g., "my_presentation.pptx").
        requests: A list of operation requests. Each request is a dictionary
                  with a single key being the operation name (e.g., "create_slide")
                  and its value being a dictionary of parameters for that operation.
                  User-defined string "object_id"s (e.g. "slide_object_id", "shape_object_id")
                  can be used to create temporary references to slides/shapes created
                  within the same batch.

    Returns:
        A dictionary containing the "presentation_id" and a list of "replies",
        one for each input request, echoing any "object_id" provided in the request.

    Example of `requests` structure:
    ```json
    [
        {
            "create_slide": {
                "layout_index": 6, // Blank layout
                "slide_object_id": "newDiagramSlide" // User-defined ID for this new slide
            }
        },
        {
            "add_shape": {
                "page_object_id": "newDiagramSlide", // Reference to the slide created above
                "shape_type_name": "RECTANGLE",
                "left_inches": 1.0, "top_inches": 1.0, "width_inches": 2.0, "height_inches": 1.0,
                "text": "Step 1: Start",
                "shape_object_id": "step1Rect" // User-defined ID for this shape
            }
        },
        {
            "add_textbox": {
                "page_object_id": "newDiagramSlide",
                "text": "Flowchart Title",
                "left_inches": 1.0, "top_inches": 0.2, "width_inches": 8.0, "height_inches": 0.5,
                "font_size_pt": 24,
                "bold": true,
                "shape_object_id": "titleBox"
            }
        },
        {
            "add_shape": {
                "page_object_id": "newDiagramSlide",
                "shape_type_name": "FLOWCHART_DECISION",
                "left_inches": 1.0, "top_inches": 3.0, "width_inches": 2.0, "height_inches": 1.5,
                "text": "Is condition met?",
                "shape_object_id": "decision1"
            }
        },
        {
            "modify_shape": {
                "page_object_id": "newDiagramSlide", // Slide containing the shape
                "shape_object_id": "step1Rect",    // ID of the shape to modify
                "text": "Step 1: Initialization",
                "fill_color_rgb": [220, 220, 255] // Light blue fill
            }
        },
        {
            "add_connector": {
                "page_object_id": "newDiagramSlide",
                "start_shape_object_id": "step1Rect", // Connect from this shape
                "end_shape_object_id": "decision1",   // To this shape
                "connector_type_name": "ELBOW",
                "shape_object_id": "connectorArrow1" // User-defined ID for the connector
            }
        },
        {
            "add_picture": {
                "page_object_id": "newDiagramSlide",
                "image_base64": "iVBORw0KGgoAAAANSUhEUgAAAAEAAAABCAQAAAC1HAwCAAAAC0lEQVR42mNkYAAAAAYAAjCB0C8AAAAASUVORK5CYII=", // Example: 1x1 black pixel PNG
                "left_inches": 5.0, "top_inches": 3.0, "width_inches": 1.0, "height_inches": 1.0,
                "shape_object_id": "tinyImage"
            }
        },
        {
            "delete_shape": {
                "page_object_id": "newDiagramSlide",
                "shape_object_id": "titleBox" // Example: delete the title box we added earlier
            }
        }
    ]
    ```
    """
    prs = _load_presentation(filename)
    object_map: Dict[str, Any] = {} # Stores user_string_id -> pptx.slide.Slide or shape_id (int)
    replies: List[Dict[str, Any]] = []

    # For atomicity: If we decide to work on a copy. For now, direct modification.
    # If any step fails, an exception will be raised, and _save_presentation won't be called.

    for req_idx, request_item in enumerate(requests):
        if not isinstance(request_item, dict) or len(request_item) != 1:
            raise ValueError(f"Request at index {req_idx} is malformed: {request_item}. Expected a dict with a single operation key.")

        operation_name, params = list(request_item.items())[0]
        if not isinstance(params, dict):
            raise ValueError(f"Parameters for operation '{operation_name}' at index {req_idx} must be a dict: {params}")

        current_reply: Dict[str, Any] = {}

        try:
            if operation_name == "create_slide":
                layout_index = params.get("layout_index", 6) # Default to blank
                slide_object_id = params.get("slide_object_id")

                if not (0 <= layout_index < len(prs.slide_layouts)):
                    raise ValueError(f"Invalid layout_index {layout_index}. Must be between 0 and {len(prs.slide_layouts) - 1}.")
                slide_layout = prs.slide_layouts[layout_index]
                new_slide = prs.slides.add_slide(slide_layout)

                if slide_object_id:
                    if not isinstance(slide_object_id, str):
                        raise ValueError(f"'{operation_name}' param 'slide_object_id' must be a string.")
                    object_map[slide_object_id] = new_slide # Store the Slide object
                    current_reply = {"create_slide": {"object_id": slide_object_id}}
                else:
                    current_reply = {"create_slide": {}} # No specific ID to echo

            elif operation_name in ["add_shape", "add_textbox"]:
                page_ref = params.get("page_object_id")
                if page_ref is None: raise ValueError(f"'{operation_name}' requires 'page_object_id'.")
                target_slide = _resolve_slide_obj(prs, page_ref, object_map, filename)

                left_in = params.get("left_inches")
                top_in = params.get("top_inches")
                width_in = params.get("width_inches")
                height_in = params.get("height_inches")
                text = params.get("text")
                shape_object_id = params.get("shape_object_id")

                if any(v is None for v in [left_in, top_in, width_in, height_in]):
                    raise ValueError(f"'{operation_name}' requires 'left_inches', 'top_inches', 'width_inches', 'height_inches'.")

                left, top = Inches(left_in), Inches(top_in)
                width, height = Inches(width_in), Inches(height_in)
                
                new_shape: Optional[BaseShape] = None
                if operation_name == "add_textbox":
                    new_shape = target_slide.shapes.add_textbox(left, top, width, height)
                    if text is not None:
                        tf = new_shape.text_frame
                        tf.clear()
                        lines = str(text).split('\\n') # Allow escaped newlines in JSON
                        p = tf.paragraphs[0]
                        p.text = lines[0]
                        for line in lines[1:]:
                            p = tf.add_paragraph()
                            p.text = line
                        tf.word_wrap = True

                    font_size_pt = params.get("font_size_pt")
                    bold = params.get("bold")
                    for p in new_shape.text_frame.paragraphs:
                        if font_size_pt is not None and font_size_pt > 0:
                            p.font.size = Pt(font_size_pt)
                        if bold is not None:
                            p.font.bold = bool(bold)
                
                elif operation_name == "add_shape":
                    shape_type_name = params.get("shape_type_name")
                    if not shape_type_name: raise ValueError("'add_shape' requires 'shape_type_name'.")
                    shape_enum = _parse_shape_type(shape_type_name)
                    new_shape = target_slide.shapes.add_shape(shape_enum, left, top, width, height)
                    if text is not None:
                        tf = new_shape.text_frame
                        tf.clear()
                        lines = str(text).split('\\n')
                        p = tf.paragraphs[0]
                        p.text = lines[0]
                        for line in lines[1:]:
                            p = tf.add_paragraph()
                            p.text = line
                        tf.word_wrap = True
                
                if new_shape and shape_object_id:
                    if not isinstance(shape_object_id, str):
                        raise ValueError(f"'{operation_name}' param 'shape_object_id' must be a string.")
                    object_map[shape_object_id] = new_shape.shape_id # Store actual int shape_id
                    current_reply = {operation_name: {"object_id": shape_object_id}}
                elif new_shape :
                     current_reply = {operation_name: {"shape_id": new_shape.shape_id}} # return actual ID if no temp was given
                else: # Should not happen if params are validated
                    current_reply = {operation_name: {}}


            elif operation_name == "add_picture":
                page_ref = params.get("page_object_id")
                if page_ref is None: raise ValueError(f"'add_picture' requires 'page_object_id'.")
                target_slide = _resolve_slide_obj(prs, page_ref, object_map, filename)

                image_base64 = params.get("image_base64")
                if not image_base64 or not isinstance(image_base64, str):
                    raise ValueError("'add_picture' requires 'image_base64' as a non-empty string.")
                
                try:
                    image_bytes = base64.b64decode(image_base64)
                except Exception as e:
                    raise ValueError(f"Invalid base64 data for image: {e}")

                image_stream = io.BytesIO(image_bytes)
                left_in = params.get("left_inches")
                top_in = params.get("top_inches")
                width_in = params.get("width_inches") # Optional
                height_in = params.get("height_inches") # Optional
                shape_object_id = params.get("shape_object_id")

                if left_in is None or top_in is None:
                     raise ValueError("'add_picture' requires 'left_inches' and 'top_inches'.")

                left, top = Inches(left_in), Inches(top_in)
                width = Inches(width_in) if width_in is not None else None
                height = Inches(height_in) if height_in is not None else None

                new_pic = target_slide.shapes.add_picture(image_stream, left, top, width=width, height=height)

                if shape_object_id:
                    if not isinstance(shape_object_id, str):
                        raise ValueError(f"'add_picture' param 'shape_object_id' must be a string.")
                    object_map[shape_object_id] = new_pic.shape_id
                    current_reply = {"add_picture": {"object_id": shape_object_id}}
                else:
                    current_reply = {"add_picture": {"shape_id": new_pic.shape_id}}


            elif operation_name == "modify_shape":
                page_ref = params.get("page_object_id") # Needed to locate the shape by its ID
                shape_ref = params.get("shape_object_id")
                if page_ref is None: raise ValueError(f"'modify_shape' requires 'page_object_id'.")
                if shape_ref is None: raise ValueError(f"'modify_shape' requires 'shape_object_id'.")

                target_slide = _resolve_slide_obj(prs, page_ref, object_map, filename)
                shape_to_modify = _resolve_shape_obj(target_slide, shape_ref, object_map)

                # Apply modifications similar to the `modify_shape` tool
                changes_made = [] # For detailed reply, if needed
                if "text" in params:
                    if shape_to_modify.has_text_frame:
                        tf = shape_to_modify.text_frame
                        tf.clear()
                        lines = str(params["text"]).split('\\n')
                        tf.text = lines[0] # First line
                        for line in lines[1:]:
                            p = tf.add_paragraph()
                            p.text = line
                        tf.word_wrap = True
                        changes_made.append("text")
                if "left_inches" in params: shape_to_modify.left = Inches(params["left_inches"]); changes_made.append("left")
                if "top_inches" in params: shape_to_modify.top = Inches(params["top_inches"]); changes_made.append("top")
                if "width_inches" in params: shape_to_modify.width = Inches(params["width_inches"]); changes_made.append("width")
                if "height_inches" in params: shape_to_modify.height = Inches(params["height_inches"]); changes_made.append("height")

                if shape_to_modify.has_text_frame:
                    if "font_size_pt" in params and params["font_size_pt"] > 0:
                        for p in shape_to_modify.text_frame.paragraphs: p.font.size = Pt(params["font_size_pt"])
                        changes_made.append("font_size")
                    if "bold" in params:
                        for p in shape_to_modify.text_frame.paragraphs: p.font.bold = bool(params["bold"])
                        changes_made.append("font_bold")
                
                if "fill_color_rgb" in params:
                    rgb = params["fill_color_rgb"]
                    if isinstance(rgb, list) and len(rgb) == 3:
                        shape_to_modify.fill.solid()
                        shape_to_modify.fill.fore_color.rgb = RGBColor(*rgb)
                        changes_made.append("fill_color")
                    else: print(f"Warning: Invalid fill_color_rgb {rgb} for batch op modify_shape.")

                if "line_color_rgb" in params or "line_width_pt" in params:
                    line = shape_to_modify.line
                    if "line_color_rgb" in params:
                        rgb_line = params["line_color_rgb"]
                        if isinstance(rgb_line, list) and len(rgb_line) == 3:
                            line.color.rgb = RGBColor(*rgb_line)
                            changes_made.append("line_color")
                        else: print(f"Warning: Invalid line_color_rgb {rgb_line} for batch op modify_shape.")
                    if "line_width_pt" in params and params["line_width_pt"] > 0:
                        line.width = Pt(params["line_width_pt"])
                        changes_made.append("line_width")
                
                # Echo object_id if it was a string reference
                obj_id_reply_val = shape_ref if isinstance(shape_ref, str) else shape_to_modify.shape_id
                current_reply = {"modify_shape": {"object_id": obj_id_reply_val, "changes": changes_made}}


            elif operation_name == "add_connector":
                page_ref = params.get("page_object_id")
                start_shape_ref = params.get("start_shape_object_id")
                end_shape_ref = params.get("end_shape_object_id")
                if None in [page_ref, start_shape_ref, end_shape_ref]:
                    raise ValueError("'add_connector' requires 'page_object_id', 'start_shape_object_id', and 'end_shape_object_id'.")

                target_slide = _resolve_slide_obj(prs, page_ref, object_map, filename)
                start_shape = _resolve_shape_obj(target_slide, start_shape_ref, object_map)
                end_shape = _resolve_shape_obj(target_slide, end_shape_ref, object_map)

                connector_type_name = params.get("connector_type_name", "ELBOW")
                start_conn_idx = params.get("start_connection_point_idx", 3)
                end_conn_idx = params.get("end_connection_point_idx", 1)
                shape_object_id = params.get("shape_object_id") # For the connector itself

                try:
                    connector_enum = getattr(MSO_CONNECTOR, connector_type_name.upper())
                except AttributeError:
                    raise ValueError(f"Unknown connector type '{connector_type_name}'. Try: STRAIGHT, ELBOW, CURVE.")

                # Initial arbitrary placement for connector
                connector = target_slide.shapes.add_connector(connector_enum, Inches(0), Inches(0), Inches(1), Inches(1))
                
                try:
                    connector.begin_connect(start_shape, start_conn_idx)
                except Exception as e: # Catch specific pptx exceptions if known, e.g. IndexError for bad conn_idx
                    print(f"Warning: Batch 'add_connector': Could not connect start to shape {start_shape.shape_id} (ref '{start_shape_ref}') at point {start_conn_idx}. Error: {e}. Attempting center.")
                    try: connector.begin_connect(start_shape, 0) # Fallback to center
                    except Exception as e2: print(f"Batch 'add_connector': Fallback start connection also failed: {e2}")
                
                try:
                    connector.end_connect(end_shape, end_conn_idx)
                except Exception as e:
                    print(f"Warning: Batch 'add_connector': Could not connect end to shape {end_shape.shape_id} (ref '{end_shape_ref}') at point {end_conn_idx}. Error: {e}. Attempting center.")
                    try: connector.end_connect(end_shape, 0) # Fallback to center
                    except Exception as e2: print(f"Batch 'add_connector': Fallback end connection also failed: {e2}")

                if shape_object_id:
                    if not isinstance(shape_object_id, str):
                        raise ValueError(f"'add_connector' param 'shape_object_id' must be a string.")
                    object_map[shape_object_id] = connector.shape_id
                    current_reply = {"add_connector": {"object_id": shape_object_id}}
                else:
                    current_reply = {"add_connector": {"shape_id": connector.shape_id}}
            
            elif operation_name == "delete_shape":
                page_ref = params.get("page_object_id")
                shape_ref = params.get("shape_object_id")
                if page_ref is None: raise ValueError(f"'delete_shape' requires 'page_object_id'.")
                if shape_ref is None: raise ValueError(f"'delete_shape' requires 'shape_object_id'.")

                target_slide = _resolve_slide_obj(prs, page_ref, object_map, filename)
                shape_to_delete = _resolve_shape_obj(target_slide, shape_ref, object_map)

                sp_el = shape_to_delete._element
                sp_el.getparent().remove(sp_el)
                
                # If shape_ref was a string ID, remove from object_map as it's gone
                if isinstance(shape_ref, str) and shape_ref in object_map:
                    del object_map[shape_ref]
                
                current_reply = {"delete_shape": {"object_id": shape_ref if isinstance(shape_ref, str) else shape_to_delete.shape_id}}

            else:
                raise ValueError(f"Unsupported operation in batch: '{operation_name}'")

        except Exception as e:
            # Attach context to the error
            err_msg = f"Error processing batch request item {req_idx} ('{operation_name}'): {e}"
            # Optionally include params in error for debugging, be careful with sensitive data
            # err_msg += f" Parameters: {params}"
            print(f"ERROR in batch_update: {err_msg}", file=sys.stderr)
            raise ValueError(err_msg) from e # Re-raise to abort batch

        replies.append(current_reply if current_reply else {operation_name: {}}) # Ensure a reply for each request

    _save_presentation(prs, filename)
    return {"presentation_id": filename, "replies": replies}

@mcp.custom_route("/ok", methods=["GET"])
async def health_check(request):
    return JSONResponse({"ok": True})


# --- Running the Server ---
if __name__ == "__main__":
    try:
        soffice_path = _find_soffice()
        print(f"✅ Found LibreOffice executable: {soffice_path}")
    except Exception as e:
        print("❌ Image rendering (`/image.png` resource) requires LibreOffice.")
        print(f"❌ {e}")
        print("   Image rendering resource will likely fail.")
    print("-" * 30)
    print(f"🚀 Starting FastMCP server for PowerPoint generation on {HOST}:{PORT}...")
    mcp.run(transport="streamable-http", host=HOST, port=PORT)
